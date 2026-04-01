"""
Enhanced Export Service - Professional DOCX and PDF exports with embedded charts
NO EMOJI in exports - Professional formatting only
"""
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from io import BytesIO
from datetime import datetime
import logging
from typing import Dict, List, Optional
import base64
from pathlib import Path
import tempfile

try:
    import plotly.io as pio
    PLOTLY_AVAILABLE = True
except ImportError:
    PLOTLY_AVAILABLE = False
    logging.warning("Plotly not available for chart exports")

logger = logging.getLogger(__name__)


class EnhancedExportService:
    """Service for generating professional DOCX and PDF exports"""
    
    def __init__(self):
        self._export_seen_blocks = set()
        # Configure plotly for static image export if available
        if PLOTLY_AVAILABLE:
            try:
                pio.kaleido.scope.mathjax = None
            except:
                pass
    
    def generate_comprehensive_docx(
        self,
        customer_name: str,
        profile_data: Dict,
        customer_data: Dict,
        market_intel: Dict = None,
        projects: List[Dict] = None,
        financial_data: Dict = None,
        charts: Dict = None,
        crm_history: Dict = None,       # from historical_service.get_yearly_performance()
        ib_data: Dict = None,           # from historical_service.get_ib_summary()
    ) -> BytesIO:
        """
        Generate comprehensive 10+ page customer analysis DOCX.
        """
        self._export_seen_blocks = set()
        doc = Document()
        self._setup_document_styles(doc)
        self._add_title_page(doc, customer_name)
        self._add_table_of_contents(doc)

        # Auto-populate market_intel from profile_data if not explicitly provided
        if not market_intel and profile_data:
            mi = profile_data.get('market_intelligence', {})
            if mi:
                market_intel = mi

        # Pg 1-2: Customer Profile (Restructured)
        self._add_customer_profile_section(doc, profile_data, customer_data, charts)

        # Pg 3: Priority Ranking Analysis
        self._add_priority_ranking_section(doc, profile_data)

        # Pg 4: Project History & Sales Relationship
        self._add_sales_relationship_section(doc, profile_data)

        # Pg 5: Customer Interactions
        self._add_customer_interactions_section(doc, profile_data)

        # Pg 6: CRM Historical Performance
        self._add_historical_crm_section(doc, crm_history, customer_name)

        # Pg 7: Deep Dive Analytics
        if customer_data:
            self._add_deep_dive_section(doc, customer_data, charts)

        # Pg 8: Market Intelligence
        if market_intel:
            self._add_market_intelligence_section(doc, market_intel)

        # Pg 9: Country-Level External Intelligence
        self._add_country_intelligence_section(doc, profile_data)

        # Pg 10: Installed Base Summary (Axel's IB list)
        self._add_installed_base_section(doc, ib_data, customer_name)

        # Pg 11: Project Analysis
        if projects:
            self._add_project_section(doc, projects, charts)

        # Pg 12: Metallurgical & Technical Insights
        self._add_metallurgical_section(doc, profile_data)

        # Pg 13: Strategic Sales Pitch
        self._add_sales_strategy_section(doc, profile_data)

        # Pg 14: Financial Analysis
        if financial_data:
            self._add_financial_section(doc, financial_data, charts)

        # Pg 14: Recent News & Developments
        self._add_recent_news_section(doc, profile_data)

        # Pg 15: References
        self._add_references_section(doc, profile_data)

        self._add_footer(doc)
        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer
    
    def _setup_document_styles(self, doc: Document):
        """Setup professional document styles"""
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = Pt(11)

        # Heading 1
        h1 = doc.styles['Heading 1']
        h1.font.size = Pt(18)
        h1.font.bold = True
        h1.font.color.rgb = RGBColor(0, 51, 102)

        # Heading 2
        h2 = doc.styles['Heading 2']
        h2.font.size = Pt(14)
        h2.font.bold = True
        h2.font.color.rgb = RGBColor(0, 102, 204)

        # Heading 3
        h3 = doc.styles['Heading 3']
        h3.font.size = Pt(12)
        h3.font.italic = True
        h3.font.color.rgb = RGBColor(50, 50, 50)

        # List bullet style
        if 'List Bullet' in doc.styles:
            bullet = doc.styles['List Bullet']
            bullet.paragraph_format.left_indent = Inches(0.25)
        
        # Add custom list bullet 2
        if 'List Bullet 2' not in doc.styles:
            doc.styles.add_style('List Bullet 2', WD_STYLE_TYPE.PARAGRAPH)
            normal = doc.styles['List Bullet 2']
            normal.font.size = Pt(11)
    
    def _add_title_page(self, doc: Document, customer_name: str):
        """Add professional title page with logo"""
        # Add SMS group logo
        logo_path = Path(__file__).resolve().parent.parent.parent / "assets" / "logo.png"
        if logo_path.exists():
            doc.add_picture(str(logo_path), width=Inches(1.5))
            last_para = doc.paragraphs[-1]
            last_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        
        doc.add_paragraph()
        doc.add_paragraph()
        doc.add_paragraph()

        # Title
        title = doc.add_paragraph()
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = title.add_run('Customer Analysis Report')
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 51, 102)
        
        doc.add_paragraph()  # Spacing
        
        # Customer name
        customer_para = doc.add_paragraph()
        customer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = customer_para.add_run(customer_name)
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0, 102, 204)
        
        doc.add_paragraph()
        doc.add_paragraph()
        
        # Generation date
        date_para = doc.add_paragraph()
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = date_para.add_run(f"Generated: {datetime.now().strftime('%B %d, %Y')}")
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(100, 100, 100)
        
        doc.add_page_break()
    
    def _add_table_of_contents(self, doc: Document):
        """Add table of contents"""
        doc.add_heading('Table of Contents', level=1)

        toc_items = [
            '1. Customer Profile',
            '2. Priority Ranking Analysis',
            '3. Project History & Sales Relationship',
            '4. Customer Interactions',
            '5. CRM Historical Performance',
            '6. Deep Dive Analytics',
            '7. Market Intelligence',
            '8. Country-Level External Intelligence',
            '9. Installed Base — Axel\'s IB List',
            '10. Project Analysis (Detailed)',
            '11. Metallurgical & Technical Insights',
            '12. Strategic Sales Pitch',
            '13. Financial Analysis',
            '14. Recent News & Developments',
            '15. References',
        ]
        for item in toc_items:
            doc.add_paragraph(item, style='List Bullet')
        doc.add_page_break()

    def _normalize_export_text(self, value: str) -> str:
        text = str(value or '').strip().lower()
        text = ' '.join(text.split())
        return text

    def _add_unique_paragraph(self, doc: Document, value: str) -> bool:
        normalized = self._normalize_export_text(value)
        if not normalized or normalized in {'n/a', 'none', 'null', '{}', '[]'}:
            return False
        if normalized in self._export_seen_blocks:
            return False
        self._export_seen_blocks.add(normalized)
        doc.add_paragraph(str(value))
        return True

    def _add_unique_heading_paragraph(self, doc: Document, heading: str, value: str, level: int = 2) -> bool:
        if self._add_unique_paragraph(doc, value):
            last_para = doc.paragraphs[-1]
            paragraph_xml = last_para._p
            body = paragraph_xml.getparent()
            body.remove(paragraph_xml)
            doc.add_heading(heading, level=level)
            doc._body._element.append(paragraph_xml)
            return True
        return False
    
    def _add_customer_profile_section(self, doc: Document, profile_data: Dict, customer_data: Dict, charts: Dict):
        """Add restructured customer profile section"""
        doc.add_heading('1. Customer Profile', level=1)
        
        # 1.1 Locations of customer on the map
        doc.add_heading('1.1 Locations on Map', level=2)
        if charts and 'Locations Map' in charts:
            self._add_chart_to_doc(doc, charts['Locations Map'], 'Locations Map')
        else:
            doc.add_paragraph('No location map available.')

        # 1.2 Basic Information
        doc.add_heading('1.2 Basic Information', level=2)
        basic_data = profile_data.get('basic_data', {})
        table = doc.add_table(rows=0, cols=2)
        table.style = 'Light Grid Accent 1'
        
        fields = [
            ('Company Name', basic_data.get('name', 'N/A')),
            ('Headquarters', basic_data.get('hq_address', 'N/A')),
            ('CEO', basic_data.get('ceo', 'N/A')),
            ('Ownership', basic_data.get('owner', 'N/A')),
            ('Employees (FTE)', basic_data.get('fte', 'N/A')),
            ('Industry', basic_data.get('company_focus', 'N/A')),
            ('Financial Status', basic_data.get('financials', 'N/A'))
        ]
        
        for label, value in fields:
            row = table.add_row()
            row.cells[0].text = label
            row.cells[1].text = str(value)
            row.cells[0].paragraphs[0].runs[0].font.bold = True

        # Additional corporate sections requested by stakeholders
        for heading, key in [
            ('1.2.1 Corporate History & Origin', 'corporate_history'),
            ('1.2.2 Shareholding / Capital Structure', 'capital_structure'),
            ('1.2.3 Employees Breakdown', 'employee_breakdown'),
            ('1.2.4 Executive Board & Management Organization', 'executive_board'),
            ('1.2.5 Group Subsidiaries / Affiliates Overview', 'subsidiaries'),
        ]:
            val = basic_data.get(key, '')
            if val:
                self._add_unique_heading_paragraph(doc, heading, str(val), level=3)

        # 1.3 Plant Locations (filtered to customer primary country)
        locs = profile_data.get('locations', [])
        if locs:
            doc.add_heading('1.3 Plant Locations', level=2)
            filter_country = (customer_data or {}).get('crm_country', '')
            if filter_country:
                filtered_locs = [loc for loc in locs if str(loc.get('country', '')).strip().lower() == filter_country.lower()]
                filtered_locs = filtered_locs if filtered_locs else locs[:20]
            else:
                filtered_locs = locs[:20]
            loc_table = doc.add_table(rows=0, cols=4)
            loc_table.style = 'Light Grid Accent 1'
            hdr_row = loc_table.add_row()
            for j, h in enumerate(['City', 'Country', 'Products', 'Capacity (t/y)']):
                hdr_row.cells[j].text = h
                if hdr_row.cells[j].paragraphs[0].runs:
                    hdr_row.cells[j].paragraphs[0].runs[0].font.bold = True
            for loc in filtered_locs:
                data_row = loc_table.add_row()
                data_row.cells[0].text = str(loc.get('city', ''))
                data_row.cells[1].text = str(loc.get('country', ''))
                data_row.cells[2].text = str(loc.get('final_products', ''))
                data_row.cells[3].text = str(loc.get('tons_per_year', ''))

        # 1.4 Equipment distribution
        doc.add_heading('1.4 Equipment Distribution', level=2)
        if charts:
            if 'Equipment Distribution' in charts:
                self._add_chart_to_doc(doc, charts['Equipment Distribution'], 'Portfolio Mix')
            
            # Additional Fleet Statistics
            for chart_name in ['Status Distribution', 'Age Distribution', 'Capacity Profile']:
                if chart_name in charts:
                    doc.add_heading(f'Fleet Insight: {chart_name}', level=3)
                    self._add_chart_to_doc(doc, charts[chart_name], chart_name)

        # 1.5 Statistical data analysis
        doc.add_heading('1.5 Statistical Data Analysis', level=2)
        stat = profile_data.get('statistical_interpretations', {})
        if stat and stat.get('charts_explanation'):
            self._add_unique_paragraph(doc, str(stat.get('charts_explanation', '')))
        else:
            doc.add_paragraph('Detailed distribution analysis of the installed base portfolio.')

        # Company Overview (Historical part of section 1)
        if 'company_overview' in profile_data:
            doc.add_heading('Company Overview', level=2)
            overview = profile_data['company_overview']
            self._add_unique_paragraph(doc, overview.get('description', 'No description available'))
            
            if overview.get('source_url'):
                p = doc.add_paragraph()
                p.add_run('Source: ')
                p.add_run(overview['source_url']).font.italic = True
        
        doc.add_page_break()

    def _add_recent_news_section(self, doc: Document, profile_data: Dict):
        """Add Recent News & Developments section before references"""
        doc.add_heading('14. Recent News & Developments', level=1)
        if 'recent_news' in profile_data:
            news_items = profile_data['recent_news']
            for idx, news in enumerate(news_items[:10], start=1):
                # Title line as a numbered, bold heading
                p = doc.add_paragraph(style='List Number')
                title = news.get('title', 'No title')
                url = news.get('url')
                title_run = p.add_run(title)
                title_run.font.bold = True
                if url:
                    try:
                        self._add_hyperlink(p, f" [{news.get('source', 'link')}]", url)
                    except Exception:
                        pass

                published_date = news.get('published_date', '')
                source = news.get('source', '')
                meta_parts = [x for x in [published_date, source] if x]
                if meta_parts:
                    p.add_run(f"  ({' | '.join(meta_parts)})").font.color.rgb = RGBColor(120, 120, 120)

                # 3-4 line summary of the news article
                desc = str(news.get('description') or news.get('summary') or '').strip()
                if desc:
                    # Truncate to roughly 4 sentences / ~600 chars for concise summary
                    sentences = [s.strip() for s in desc.replace('\n', ' ').split('. ') if s.strip()]
                    summary = '. '.join(sentences[:4])
                    if summary and not summary.endswith('.'):
                        summary += '.'
                    if len(summary) > 650:
                        summary = summary[:647] + '...'
                    sq = doc.add_paragraph(style='Normal')
                    sq.paragraph_format.left_indent = Inches(0.3)
                    sq.paragraph_format.space_after = Pt(6)
                    sq.add_run(summary).font.size = Pt(10)
                doc.add_paragraph()  # spacing between news items
        else:
            doc.add_paragraph('No recent news available.')
        doc.add_page_break()

    def _add_hyperlink(self, paragraph, text: str, url: str):
        """Insert a clickable hyperlink run into a python-docx paragraph."""
        part = paragraph.part
        r_id = part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", is_external=True)

        hyperlink = OxmlElement('w:hyperlink')
        hyperlink.set(qn('r:id'), r_id)

        new_run = OxmlElement('w:r')
        r_pr = OxmlElement('w:rPr')

        color = OxmlElement('w:color')
        color.set(qn('w:val'), '0563C1')
        r_pr.append(color)

        underline = OxmlElement('w:u')
        underline.set(qn('w:val'), 'single')
        r_pr.append(underline)

        new_run.append(r_pr)
        text_elem = OxmlElement('w:t')
        text_elem.text = text
        new_run.append(text_elem)
        hyperlink.append(new_run)

        paragraph._p.append(hyperlink)

    def convert_docx_to_pdf(self, docx_buffer: BytesIO) -> BytesIO:
        """Convert DOCX bytes to PDF so both exports contain the same content."""
        try:
            from docx2pdf import convert
        except ImportError as exc:
            raise RuntimeError(
                "docx2pdf is required for PDF export consistency. Install with: pip install docx2pdf"
            ) from exc

        with tempfile.TemporaryDirectory(prefix='sms_export_') as tmpdir:
            tmp_path = Path(tmpdir)
            in_path = tmp_path / 'report.docx'
            out_path = tmp_path / 'report.pdf'

            in_path.write_bytes(docx_buffer.getvalue())

            try:
                convert(str(in_path), str(out_path))
            except Exception as exc:
                raise RuntimeError(
                    "DOCX-to-PDF conversion failed. Ensure Microsoft Word is installed and accessible on this machine."
                ) from exc

            if not out_path.exists():
                raise RuntimeError("PDF file was not produced by docx2pdf conversion.")

            pdf_buffer = BytesIO(out_path.read_bytes())
            pdf_buffer.seek(0)
            return pdf_buffer
    
    def _add_deep_dive_section(self, doc: Document, customer_data: Dict, charts: Dict):
        """Add deep dive analytics section"""
        doc.add_heading('6. Deep Dive Analytics', level=1)

        projects_dd   = customer_data.get('projects', [])
        installed_dd  = customer_data.get('installed_base', [])
        total_rev     = sum(p.get('value', 0) for p in projects_dd)
        active_projects = sum(1 for p in projects_dd if p.get('status') in ['Active', 'In Progress'])

        doc.add_heading('Key Performance Indicators', level=2)
        table = doc.add_table(rows=0, cols=2)
        table.style = 'Light Grid Accent 1'
        for label, val in [
            ('Total Projects',   str(len(projects_dd))),
            ('Active Projects',  str(active_projects)),
            ('Total Equipment (CRM Export)',  str(len(installed_dd))),
            ('Total Revenue',    f'EUR {total_rev:,.0f}'),
        ]:
            row = table.add_row()
            row.cells[0].text = label
            row.cells[1].text = val
            row.cells[0].paragraphs[0].runs[0].font.bold = True

        if charts and 'revenue_trend' in charts:
            self._add_chart_to_doc(doc, charts['revenue_trend'], 'Revenue Trend Chart')

        doc.add_page_break()
    
    def _add_market_intelligence_section(self, doc: Document, market_intel: Dict):
        """Add market intelligence section"""
        doc.add_heading('7. Market Intelligence', level=1)

        sections = [
            ('Financial Health',    market_intel.get('financial_health', '')),
            ('Recent Developments', market_intel.get('recent_developments', '')),
            ('Market Position',     market_intel.get('market_position', '')),
            ('Strategic Outlook',   market_intel.get('strategic_outlook', '')),
            ('Risk Assessment',     market_intel.get('risk_assessment', '')),
            ('Market Size',         market_intel.get('market_size', '')),
            ('Growth Trends',       market_intel.get('growth_trends', '')),
        ]
        any_content = False
        for title, content in sections:
            val = content
            if isinstance(val, dict):
                val = val.get('summary', '') or val.get('text', '') or str(val)
            if val:
                added = self._add_unique_heading_paragraph(doc, title, str(val), level=2)
                any_content = any_content or added

        if not any_content:
            doc.add_paragraph('Market intelligence data not available. Run profile generation to populate this section.')

        if market_intel.get('competitors'):
            doc.add_heading('Key Competitors', level=2)
            for competitor in market_intel['competitors']:
                doc.add_paragraph(competitor, style='List Bullet')

        workforce_strategy = market_intel.get('workforce_strategy', '')
        if workforce_strategy:
            self._add_unique_heading_paragraph(doc, 'Workforce Strategy', str(workforce_strategy), level=2)

        product_portfolio = market_intel.get('product_portfolio', '')
        if product_portfolio:
            self._add_unique_heading_paragraph(doc, 'Products of Steel Business Unit', str(product_portfolio), level=2)

        end_market_breakdown = market_intel.get('end_market_breakdown', '')
        if end_market_breakdown:
            self._add_unique_heading_paragraph(doc, 'Construction / Automotive-specific Market Sections', str(end_market_breakdown), level=2)

        # NOTE: sources/references are consolidated in the final References section (Section 16)
        doc.add_page_break()
    
    def _add_project_section(self, doc: Document, projects: List[Dict], charts: Dict):
        """Add project analysis section"""
        doc.add_heading('10. Project Analysis', level=1)

        doc.add_heading('Project Summary', level=2)
        doc.add_paragraph(f'Total Projects: {len(projects)}')
        active = sum(1 for p in projects if p.get('status') in ['Active', 'In Progress'])
        completed = sum(1 for p in projects if p.get('status') in ['Completed', 'Won'])
        total_val = sum(p.get('value', 0) for p in projects)
        doc.add_paragraph(f'Active Projects: {active}')
        doc.add_paragraph(f'Completed Projects: {completed}')
        doc.add_paragraph(f'Total Value: EUR {total_val:,.0f}')

        for project in projects:
            doc.add_heading(project.get('name', 'Unnamed Project'), level=3)
            details = [
                ('Status',    project.get('status', 'Unknown')),
                ('Start',     project.get('start_date', 'N/A')),
                ('End',       project.get('end_date', 'N/A')),
                ('Value',     f"EUR {project.get('value', 0):,.0f}"),
                ('Budget',    f"EUR {project.get('budget', 0):,.0f}"),
                ('Progress',  f"{project.get('progress', 0)}%"),
                ('Type',      project.get('type', 'N/A')),
            ]
            table = doc.add_table(rows=0, cols=2)
            table.style = 'Light List Accent 1'
            for label, value in details:
                row = table.add_row()
                row.cells[0].text = label
                row.cells[1].text = str(value)
            doc.add_paragraph()

        if charts and 'gantt_chart' in charts:
            self._add_chart_to_doc(doc, charts['gantt_chart'], 'Project Timeline')

        doc.add_page_break()
    
    def _add_financial_section(self, doc: Document, financial_data: Dict, charts: Dict):
        """Add financial analysis section"""
        doc.add_heading('13. Financial Analysis', level=1)

        doc.add_heading('Cost Breakdown', level=2)
        costs = financial_data.get('cost_breakdown', {})
        table = doc.add_table(rows=0, cols=2)
        table.style = 'Medium Grid 1 Accent 1'
        header = table.add_row()
        header.cells[0].text = 'Category'
        header.cells[1].text = 'Amount (EUR)'
        for cell in header.cells:
            if cell.paragraphs[0].runs:
                cell.paragraphs[0].runs[0].font.bold = True
        for category, amount in costs.items():
            row = table.add_row()
            row.cells[0].text = str(category)
            row.cells[1].text = f"EUR {float(amount):,.2f}" if str(amount).replace('.','').isdigit() else str(amount)

        if 'budget_variance' in financial_data:
            doc.add_heading('Budget Variance Analysis', level=2)
            variance = financial_data['budget_variance']
            doc.add_paragraph(f"Budgeted: EUR {variance.get('budgeted', 0):,.2f}")
            doc.add_paragraph(f"Actual:   EUR {variance.get('actual', 0):,.2f}")
            doc.add_paragraph(f"Variance: EUR {variance.get('variance', 0):,.2f} ({variance.get('variance_percent', 0):.1f}%)")
            doc.add_paragraph(f"Status: {variance.get('status', 'Unknown')}")

        if charts and 'cost_trend' in charts:
            self._add_chart_to_doc(doc, charts['cost_trend'], 'Cost Trend Analysis')
        doc.add_page_break()
    
    def _add_chart_to_doc(self, doc: Document, fig, caption: str):
        """Add plotly chart as image to document"""
        if not PLOTLY_AVAILABLE:
            doc.add_paragraph(f"[Chart: {caption}] - Chart export requires kaleido library")
            return
        
        try:
            # Convert plotly figure to image bytes
            img_bytes = pio.to_image(fig, format='png', width=800, height=400)
            img_stream = BytesIO(img_bytes)
            
            # Add to document
            doc.add_picture(img_stream, width=Inches(6))
            
            # Add caption
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(caption)
            run.font.italic = True
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(100, 100, 100)
            
        except Exception as e:
            logger.error(f"Failed to add chart {caption}: {e}")
            doc.add_paragraph(f"[Chart: {caption}] - Error generating chart image")
    
    def _add_footer(self, doc: Document):
        """Add document footer with page numbers to every section."""

        def _make_rpr():
            """Create a fresh w:rPr element (8pt, grey). Must be called per-run."""
            rpr = OxmlElement('w:rPr')
            sz = OxmlElement('w:sz')
            sz.set(qn('w:val'), '16')  # 8pt
            col = OxmlElement('w:color')
            col.set(qn('w:val'), '999999')
            rpr.append(sz)
            rpr.append(col)
            return rpr

        def _text_run(para_elem, text: str):
            r = OxmlElement('w:r')
            r.append(_make_rpr())
            t = OxmlElement('w:t')
            t.set(qn('xml:space'), 'preserve')
            t.text = text
            r.append(t)
            para_elem.append(r)

        def _field_run(para_elem, field_code: str):
            """Append a PAGE or NUMPAGES field run (3 runs: begin / instr / end)."""
            r_begin = OxmlElement('w:r')
            r_begin.append(_make_rpr())
            begin = OxmlElement('w:fldChar')
            begin.set(qn('w:fldCharType'), 'begin')
            r_begin.append(begin)

            r_instr = OxmlElement('w:r')
            r_instr.append(_make_rpr())
            instr = OxmlElement('w:instrText')
            instr.set(qn('xml:space'), 'preserve')
            instr.text = f' {field_code} '
            r_instr.append(instr)

            r_end = OxmlElement('w:r')
            r_end.append(_make_rpr())
            end_fc = OxmlElement('w:fldChar')
            end_fc.set(qn('w:fldCharType'), 'end')
            r_end.append(end_fc)

            para_elem.append(r_begin)
            para_elem.append(r_instr)
            para_elem.append(r_end)

        for section in doc.sections:
            footer = section.footer
            footer.is_linked_to_previous = False
            fp = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
            fp.clear()
            fp.alignment = WD_ALIGN_PARAGRAPH.CENTER

            p_elem = fp._p
            _text_run(p_elem, f"Customer Analysis Report | CONFIDENTIAL    |    Page ")
            _field_run(p_elem, 'PAGE')
            _text_run(p_elem, " of ")
            _field_run(p_elem, 'NUMPAGES')
            _text_run(p_elem, f"    |    Generated {datetime.now().strftime('%Y-%m-%d')}")

    # ── NEW sections (CRM History, IB, Metallurgical, Sales Strategy, Relationship, Priority, Country) ─

    def _add_priority_ranking_section(self, doc: Document, profile_data: Dict):
        """Section 2: Priority Ranking Analysis"""
        doc.add_heading('2. Priority Ranking Analysis', level=1)
        priority = profile_data.get('priority_analysis', {})
        if not priority:
            doc.add_paragraph('No priority ranking data available.')
            doc.add_page_break()
            return
            
        doc.add_heading('Ranking Overview', level=2)
        table = doc.add_table(rows=0, cols=2)
        table.style = 'Light Grid Accent 1'
        for label, val in [
            ('Priority Score', priority.get('priority_score', 'N/A')),
            ('Customer Rank', priority.get('priority_rank', 'N/A')),
        ]:
            row = table.add_row()
            row.cells[0].text = label
            row.cells[1].text = str(val)
            row.cells[0].paragraphs[0].runs[0].font.bold = True

        # company_explainer holds the deep analytical narrative
        explainer = priority.get('company_explainer', priority.get('reasoning', ''))
        if explainer:
            self._add_unique_heading_paragraph(doc, 'Company Analysis', str(explainer), level=2)
            
        self._add_unique_heading_paragraph(doc, 'Key Opportunity Drivers', str(priority.get('key_opportunity_drivers', 'N/A')), level=2)
        
        self._add_unique_heading_paragraph(doc, 'Engagement Recommendation', str(priority.get('engagement_recommendation', 'N/A')), level=2)
        
        doc.add_page_break()

    def _add_country_intelligence_section(self, doc: Document, profile_data: Dict):
        """Section 7: Country-Level External Intelligence"""
        doc.add_heading('8. Country-Level External Intelligence', level=1)
        ci = profile_data.get('country_intelligence', {})
        if not ci:
            doc.add_paragraph('No country-level intelligence available.')
            doc.add_page_break()
            return
            
        fields = [
            ('Steel Market Summary', 'steel_market_summary'),
            ('Economic Context', 'economic_context'),
            ('Trade & Tariff Context', 'trade_tariff_context'),
            ('Automotive Sector Trends', 'automotive_sector'),
            ('Investment Drivers', 'investment_drivers')
        ]
        
        for title, key in fields:
            val = ci.get(key, 'N/A')
            if val and val != 'N/A':
                self._add_unique_heading_paragraph(doc, title, str(val), level=2)
        
        doc.add_page_break()

    def _add_sales_relationship_section(self, doc: Document, profile_data: Dict):
        """Section 3: Project History & Sales Relationship"""
        doc.add_heading('3. Project History & Sales Relationship', level=1)
        history = profile_data.get('history', {}) if profile_data else {}
        context = profile_data.get('context', {}) if profile_data else {}

        doc.add_heading('Sales Relationship Summary', level=2)
        for label, key in [
            ('CRM Rating', 'crm_rating'),
            ('SMS Relationship', 'sms_relationship'),
            ('Key Contact Person', 'key_person'),
            ('Latest Visits', 'latest_visits'),
        ]:
            val = history.get(key, '')
            if val:
                doc.add_paragraph(f'{label}: {val}')

        # Existing SMS installations
        sms_hist = history.get('sms_delivery_history', '')
        if sms_hist:
            self._add_unique_heading_paragraph(doc, 'Existing Facilities of SMS Group (History)', str(sms_hist), level=2)

        doc.add_heading('Project Track Record', level=2)
        for label, key in [
            ('Latest Projects', 'latest_projects'),
            ('Realized Projects', 'realized_projects'),
        ]:
            val = history.get(key, '')
            if val:
                self._add_unique_heading_paragraph(doc, label, str(val), level=3)

        # New project sections from Thomas's feedback
        for section_title, history_key in [
            ('Current Projects (Active Opportunities)', 'current_projects_detail'),
            ('Projects Under Execution', 'projects_under_execution'),
            ('Lost Projects', 'lost_projects'),
        ]:
            val = history.get(history_key, '')
            if val:
                self._add_unique_heading_paragraph(doc, section_title, str(val), level=2)

        announced = (profile_data.get('market_intelligence') or {}).get('announced_investments', '')
        if announced:
            self._add_unique_heading_paragraph(doc, 'Announced Investments (News)', str(announced), level=2)

        doc.add_heading('Market Context', level=2)
        end_cust = context.get('end_customer', '')
        mkt_pos = context.get('market_position', '')
        if end_cust:
            doc.add_paragraph(f'End Customer / Supply Chain: {end_cust}')
        if mkt_pos:
            doc.add_paragraph(f'Market Position: {mkt_pos}')

        doc.add_page_break()

    def _add_customer_interactions_section(self, doc: Document, profile_data: Dict):
        """Section 4: Customer interaction timeline."""
        doc.add_heading('4. Customer Interactions', level=1)
        summary = (profile_data or {}).get('customer_interaction_summary', {}) or {}
        interactions = (profile_data or {}).get('customer_interactions', []) or []

        if summary:
            doc.add_heading('Interaction Summary', level=2)
            table = doc.add_table(rows=0, cols=2)
            table.style = 'Light Grid Accent 1'
            for label, value in [
                ('Total Interactions', summary.get('total_interactions', 0)),
                ('Last Contact Date', summary.get('last_contact_date', 'N/A')),
                ('Last Contact Location', summary.get('last_contact_location', 'N/A')),
                ('Last Contact By', summary.get('last_contact_owner', 'N/A')),
                ('Last Contact Subject', summary.get('last_contact_subject', 'N/A')),
                ('Main Channels', ', '.join(summary.get('top_channels', [])) or 'N/A'),
                ('Main SMS Contacts', ', '.join(summary.get('top_contacts', [])) or 'N/A'),
            ]:
                row = table.add_row()
                row.cells[0].text = str(label)
                row.cells[1].text = str(value)
                if row.cells[0].paragraphs[0].runs:
                    row.cells[0].paragraphs[0].runs[0].font.bold = True

        if interactions:
            doc.add_heading('Recent Interaction Timeline', level=2)
            table = doc.add_table(rows=0, cols=5)
            table.style = 'Light Grid Accent 1'
            hdr = table.add_row()
            for idx, heading in enumerate(['Date', 'Channel / Location', 'Responsible', 'Account', 'Subject']):
                hdr.cells[idx].text = heading
                if hdr.cells[idx].paragraphs[0].runs:
                    hdr.cells[idx].paragraphs[0].runs[0].font.bold = True
            for item in interactions[:12]:
                row = table.add_row()
                start_date = str(item.get('start_dt', '') or '')[:10]
                channel_bits = [
                    str(item.get('distribution_channel', '') or '').strip(),
                    str(item.get('meeting_location', '') or '').strip(),
                ]
                row.cells[0].text = start_date
                row.cells[1].text = ' | '.join(bit for bit in channel_bits if bit)
                row.cells[2].text = str(item.get('employee_responsible', '') or '')
                row.cells[3].text = str(item.get('account', '') or '')
                row.cells[4].text = str(item.get('subject', '') or '')
        else:
            doc.add_paragraph('No recent customer interactions available in the SAP Sales Cloud visit export.')

        doc.add_page_break()

    def _add_historical_crm_section(self, doc: Document, crm_history: Optional[Dict], customer_name: str):
        """Section 5: CRM historical performance and project overview."""
        doc.add_heading('5. CRM Historical Performance', level=1)
        doc.add_heading('Annual Pipeline & Win Rate (Axel\'s CRM Export)', level=2)

        if not crm_history:
            doc.add_paragraph('No CRM historical data available for this customer.')
            doc.add_page_break()
            return

        metrics = crm_history.get('metrics', {}) if isinstance(crm_history.get('metrics', {}), dict) else {}
        overview = doc.add_table(rows=0, cols=2)
        overview.style = 'Light Grid Accent 1'
        for label, value in [
            ('Total Projects', metrics.get('n_projects', 'N/A')),
            ('Total Won Value (EUR)', f"EUR {float(metrics.get('total_won_value', 0) or 0):,.0f}"),
            ('Overall Win Rate (%)', f"{float(metrics.get('win_rate', 0) or 0):.1f}%"),
            ('Years Covered', metrics.get('time_span', 'N/A')),
        ]:
            row = overview.add_row()
            row.cells[0].text = str(label)
            row.cells[1].text = str(value)
            if row.cells[0].paragraphs[0].runs:
                row.cells[0].paragraphs[0].runs[0].font.bold = True

        yearly_df = crm_history.get('yearly_df')
        if yearly_df is not None and not yearly_df.empty:
            doc.add_heading('Year-by-Year Summary', level=3)
            rows = [['Year', 'Projects', 'Total Pipeline (EUR)', 'Won Value (EUR)', 'Win Rate %']]
            for _, r in yearly_df.iterrows():
                raw_year = r.get('Year', '')
                try:
                    year_str = str(int(float(raw_year))) if raw_year not in ('', None) else ''
                except (ValueError, TypeError):
                    year_str = str(raw_year)
                rows.append([
                    year_str,
                    str(int(r.get('Projects', 0) or 0)),
                    f"EUR {float(r.get('Total Value (EUR)', 0) or 0):,.0f}",
                    f"EUR {float(r.get('Won Value (EUR)', 0) or 0):,.0f}",
                    f"{float(r.get('Win Rate %', 0) or 0):.1f}%",
                ])
            table = doc.add_table(rows=0, cols=5)
            table.style = 'Medium Grid 1 Accent 1'
            for i, row_data in enumerate(rows):
                row = table.add_row()
                for j, cell_value in enumerate(row_data):
                    row.cells[j].text = str(cell_value)
                    if i == 0 and row.cells[j].paragraphs[0].runs:
                        row.cells[j].paragraphs[0].runs[0].font.bold = True

        raw_projects = crm_history.get('raw_projects')
        if raw_projects is not None and not raw_projects.empty:
            active_mask = raw_projects['_status'].astype(str).str.lower().str.contains('active|progress|execution|negotiation|offer|proposal|budget', na=False)
            active_projects = raw_projects[active_mask].head(12)
            if not active_projects.empty:
                doc.add_heading('Current Projects / Projects Under Execution', level=3)
                show_cols = [c for c in ['account_name', 'customer_project', '_status', '_value', '_year', 'sp_coe'] if c in active_projects.columns]
                if show_cols:
                    table = doc.add_table(rows=0, cols=len(show_cols))
                    table.style = 'Light Grid Accent 1'
                    hdr = table.add_row()
                    for idx, col in enumerate(show_cols):
                        hdr.cells[idx].text = col.replace('_', ' ').title()
                        if hdr.cells[idx].paragraphs[0].runs:
                            hdr.cells[idx].paragraphs[0].runs[0].font.bold = True
                    for _, dr in active_projects.iterrows():
                        row = table.add_row()
                        for idx, col in enumerate(show_cols):
                            val = dr.get(col, '')
                            row.cells[idx].text = f"EUR {float(val):,.0f}" if col == '_value' and str(val) not in ('', 'nan') else str(val) if val is not None else ''

        won_list = crm_history.get('won_list')
        if won_list is not None and not won_list.empty:
            doc.add_heading('Existing SMS Facilities / Won Projects History', level=3)
            show_cols = [c for c in [
                'account_name', 'codeword_sales', 'customer_project',
                'cp_expected_value_eur', '_year', 'account_country', 'sp_coe',
            ] if c in won_list.columns]
            if show_cols:
                won_sub = won_list[show_cols].head(20)
                hdr = [c.replace('_', ' ').title() for c in show_cols]
                table = doc.add_table(rows=0, cols=len(show_cols))
                table.style = 'Light Grid Accent 1'
                hdr_row = table.add_row()
                for j, h in enumerate(hdr):
                    hdr_row.cells[j].text = h
                    if hdr_row.cells[j].paragraphs[0].runs:
                        hdr_row.cells[j].paragraphs[0].runs[0].font.bold = True
                for _, dr in won_sub.iterrows():
                    drow = table.add_row()
                    for j, col in enumerate(show_cols):
                        val = dr.get(col, '')
                        drow.cells[j].text = f"EUR {float(val):,.0f}" if col == 'cp_expected_value_eur' and str(val) not in ('', 'nan') else str(val) if val is not None else ''

        lost_list = crm_history.get('lost_list')
        if lost_list is not None and not lost_list.empty:
            doc.add_heading('Lost / Non-Won Projects', level=3)
            show_cols = [c for c in ['account_name', 'customer_project', '_status', '_value', '_year'] if c in lost_list.columns]
            if show_cols:
                subset = lost_list[show_cols].head(12)
                table = doc.add_table(rows=0, cols=len(show_cols))
                table.style = 'Light Grid Accent 1'
                hdr = table.add_row()
                for idx, col in enumerate(show_cols):
                    hdr.cells[idx].text = col.replace('_', ' ').title()
                    if hdr.cells[idx].paragraphs[0].runs:
                        hdr.cells[idx].paragraphs[0].runs[0].font.bold = True
                for _, dr in subset.iterrows():
                    row = table.add_row()
                    for idx, col in enumerate(show_cols):
                        val = dr.get(col, '')
                        row.cells[idx].text = f"EUR {float(val):,.0f}" if col == '_value' and str(val) not in ('', 'nan') else str(val) if val is not None else ''

        doc.add_page_break()

    def _add_installed_base_section(self, doc: Document, ib_data: Optional[Dict], customer_name: str):
        """Section 9: Structured installed-base summary."""
        doc.add_heading('9. Installed Base (Axel\'s IB List)', level=1)

        if not ib_data or ib_data.get('n_units', 0) == 0:
            doc.add_paragraph('No installed base records found for this customer in the IB list.')
            doc.add_page_break()
            return

        doc.add_heading('IB Summary', level=2)
        table = doc.add_table(rows=0, cols=2)
        table.style = 'Light Grid Accent 1'
        for label, val in [
            ('Equipment Units (IB List)', str(ib_data.get('n_units', 0))),
            ('Average Age (years)', str(ib_data.get('avg_age', 'N/A'))),
            ('Equipment Types', ', '.join(str(t) for t in ib_data.get('equipment_types', [])) or 'N/A'),
            ('Countries / Regions', ', '.join(str(c) for c in ib_data.get('countries', [])) or 'N/A'),
        ]:
            row = table.add_row()
            row.cells[0].text = label
            row.cells[1].text = str(val)
            if row.cells[0].paragraphs[0].runs:
                row.cells[0].paragraphs[0].runs[0].font.bold = True

        df_ib = ib_data.get('df')
        if df_ib is not None and not df_ib.empty:
            doc.add_heading('Structured Installed Base Overview (top 30)', level=2)
            company_col = next((c for c in ['Company', 'Parent Company', 'ib_customer', 'account_name', 'customer'] if c in df_ib.columns), None)
            site_col = ib_data.get('city_col') or next((c for c in ['City', 'Technical Location', 'ib_city', 'site_name', 'city'] if c in df_ib.columns), None)
            equipment_col = ib_data.get('prod_col') or next((c for c in ['Type of Plant', 'Installed Base', 'ib_machine', 'ib_description', 'ib_product', 'equipment'] if c in df_ib.columns), None)
            capacity_col = next((c for c in ['Nominat Capacity [t/y]', 'Nominal Capacity', 'capacity', 'capacity_internal'] if c in df_ib.columns), None)
            year_col = ib_data.get('year_col') or next((c for c in ['Year of Start Up', 'Date Start-up', 'ib_startup', 'start_year', 'installation_year', '_year'] if c in df_ib.columns), None)

            display_rows = []
            subset = df_ib.head(30)
            for _, dr in subset.iterrows():
                year_value = dr.get(year_col, '') if year_col else ''
                age_value = dr.get('_age', '')
                if str(year_value) not in ('', 'nan', 'NaT') and str(age_value) not in ('', 'nan'):
                    year_age = f"{year_value} / {int(float(age_value))}y"
                elif str(year_value) not in ('', 'nan', 'NaT'):
                    year_age = str(year_value)
                else:
                    year_age = ''
                display_rows.append([
                    str(dr.get(company_col, customer_name)) if company_col else customer_name,
                    str(dr.get(site_col, '')) if site_col else '',
                    str(dr.get(equipment_col, '')) if equipment_col else '',
                    str(dr.get(capacity_col, '')) if capacity_col else '',
                    year_age,
                ])

            table = doc.add_table(rows=0, cols=5)
            table.style = 'Light List'
            hdr = table.add_row()
            for idx, heading in enumerate(['Company', 'Site', 'Equipment', 'Capacity', 'Year of Start-Up / Age']):
                hdr.cells[idx].text = heading
                if hdr.cells[idx].paragraphs[0].runs:
                    hdr.cells[idx].paragraphs[0].runs[0].font.bold = True
            for display_row in display_rows:
                row = table.add_row()
                for idx, value in enumerate(display_row):
                    row.cells[idx].text = '' if value in ('nan', 'None') else str(value)

        doc.add_page_break()

    def _add_metallurgical_section(self, doc: Document, profile_data: Dict):
        """Section 11: Metallurgical & Technical Insights"""
        doc.add_heading('11. Metallurgical & Technical Insights', level=1)
        meta = (profile_data or {}).get('metallurgical_insights', {})
        if not meta:
            doc.add_paragraph('No metallurgical insights available. Generate a profile first.')
            doc.add_page_break()
            return
        for label, key in [
            ('Process Efficiency', 'process_efficiency'),
            ('Carbon Footprint / Green Steel', 'carbon_footprint_strategy'),
            ('Modernization Potential', 'modernization_potential'),
            ('Technical Bottlenecks', 'technical_bottlenecks'),
        ]:
            val = meta.get(key, '')
            if val:
                self._add_unique_heading_paragraph(doc, label, str(val), level=2)
        doc.add_page_break()

    def _add_sales_strategy_section(self, doc: Document, profile_data: Dict):
        """Section 12: Strategic Sales Pitch"""
        doc.add_heading('12. Strategic Sales Pitch', level=1)
        strat = (profile_data or {}).get('sales_strategy', {})
        market = (profile_data or {}).get('market_intelligence', {})
        if not strat:
            doc.add_paragraph('No sales strategy available. Generate a profile first.')
            doc.add_page_break()
            return

        # Added per Thomas feedback: explicit commercial structure and product context
        if strat.get('sms_commercial_structure'):
            self._add_unique_heading_paragraph(doc, 'Commercial Sales SMS group - Sales Structure', str(strat.get('sms_commercial_structure')), level=2)

        if strat.get('buying_center_map'):
            self._add_unique_heading_paragraph(doc, 'Buying Center Map', str(strat.get('buying_center_map')), level=2)

        if market.get('product_portfolio'):
            self._add_unique_heading_paragraph(doc, 'Products of Steel Business Unit', str(market.get('product_portfolio')), level=2)

        if market.get('end_market_breakdown'):
            self._add_unique_heading_paragraph(doc, 'Construction / Automotive-specific Market Sections', str(market.get('end_market_breakdown')), level=2)
        for label, key in [
            ('Recommended Portfolio', 'recommended_portfolio'),
            ('Value Proposition', 'value_proposition'),
            ('Competitive Landscape', 'competitive_landscape'),
            ('Suggested Next Steps', 'suggested_next_steps'),
        ]:
            val = strat.get(key, '')
            if val:
                self._add_unique_heading_paragraph(doc, label, str(val), level=2)

        if strat.get('compliance_guidance'):
            self._add_unique_heading_paragraph(doc, 'Compliance Implications For SMS', str(strat.get('compliance_guidance')), level=2)
        if not any(strat.get(k) for k in ('recommended_portfolio', 'value_proposition', 'competitive_landscape', 'suggested_next_steps')):
            raw = str(strat)
            if raw and raw != '{}':
                doc.add_paragraph(raw)
        doc.add_page_break()
    
    def _add_statistical_charts_section(self, doc: Document, profile_data: Dict, charts: Dict):
        """Section 14: Statistical Graphs"""
        doc.add_heading('14. Statistical Graphs & Visualizations', level=1)
        
        stat = (profile_data or {}).get('statistical_interpretations', {})
        if stat and stat.get('charts_explanation'):
            doc.add_heading('Statistical Data Analysis', level=2)
            doc.add_paragraph(str(stat.get('charts_explanation', '')))
        
        for chart_name, fig in charts.items():
            doc.add_heading(chart_name, level=2)
            self._add_chart_to_doc(doc, fig, chart_name)
        
        doc.add_page_break()
        
    def _add_references_section(self, doc: Document, profile_data: Dict):
        """Section 15: References"""
        doc.add_heading('15. References', level=1)
        refs = (profile_data or {}).get('references', [])
        
        if not refs:
            doc.add_paragraph('No specific external references cited.')
        else:
            for ref in refs:
                doc.add_paragraph(str(ref), style='List Bullet')
        doc.add_page_break()

    def generate_filename(self, customer_name: str, extension: str) -> str:
        """Generate standardized filename with date and time"""
        safe_name = "".join(c for c in customer_name if c.isalnum() or c in (' ', '-', '_')).strip()
        safe_name = safe_name.replace(' ', '_')
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        return f"{safe_name}_{timestamp}.{extension}"

    def generate_comprehensive_pptx(
        self,
        customer_name: str,
        profile_data: Dict,
        customer_data: Dict,
        crm_history: Optional[Dict] = None,
        ai_slide_outline: Optional[List[Dict]] = None,
    ) -> BytesIO:
        """Generate customer profile PPTX deck with executive summary and action slides."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError as exc:
            raise RuntimeError("python-pptx is required for PPTX export. Install with: pip install python-pptx") from exc

        prs = Presentation()

        # ── Helper: add a content slide ────────────────────────────────────
        def _add_content_slide(title: str, bullets: list, layout_idx: int = 1):
            sl = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            sl.shapes.title.text = title
            tf = sl.shapes.placeholders[1].text_frame
            tf.clear()
            tf.word_wrap = True
            for i, bullet in enumerate(bullets[:7]):
                bp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                bp.text = str(bullet)[:380]
                bp.level = 0
                if bp.runs:
                    bp.runs[0].font.size = Pt(16)
            return sl

        def _add_table_slide(title: str, header: list, rows: list, layout_idx: int = 5):
            sl = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            sl.shapes.title.text = title
            n_cols = len(header)
            n_rows = min(len(rows) + 1, 15)
            tbl = sl.shapes.add_table(n_rows, n_cols, Inches(0.4), Inches(1.3), Inches(12.3), Inches(5.0)).table
            for j, h in enumerate(header):
                tbl.cell(0, j).text = str(h)
            for i, row in enumerate(rows[:n_rows - 1], start=1):
                for j, cell_val in enumerate(row[:n_cols]):
                    tbl.cell(i, j).text = str(cell_val)[:120]
            return sl

        def _truncate(text: str, chars: int = 380) -> str:
            s = str(text or '').strip()
            return (s[:chars] + '…') if len(s) > chars else s

        # ── SLIDE 1: Title ──────────────────────────────────────────────────
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = f"Customer Analysis Report"
        subtitle = slide.placeholders[1]
        subtitle.text = f"{customer_name}\nSMS Group | Generated {datetime.now().strftime('%B %d, %Y')}"

        # ── SLIDE 2: Executive Snapshot ──────────────────────────────────────
        basic = profile_data.get('basic_data', {})
        pa    = profile_data.get('priority_analysis', {})
        hist  = profile_data.get('history', {})
        _add_content_slide("Executive Snapshot", [
            f"Priority Score: {pa.get('priority_score', 'N/A')}  |  Rank: {pa.get('priority_rank', 'N/A')}",
            f"HQ: {basic.get('hq_address', 'N/A')}",
            f"Employees: {basic.get('fte', 'N/A')}  |  Owner: {basic.get('owner', 'N/A')}",
            f"Industry Focus: {basic.get('company_focus', 'N/A')}",
            f"Financial Status: {basic.get('financials', 'N/A')}",
            f"SMS Relationship: {hist.get('sms_relationship', 'N/A')}",
            f"Total Won (EUR): {hist.get('total_won_value_eur', 'N/A')}  |  Win Rate: {hist.get('win_rate_pct', 'N/A')}%",
        ])

        # ── SLIDE 3: Priority Analysis ───────────────────────────────────────
        company_explainer = pa.get('company_explainer') or pa.get('reasoning', '')
        _add_content_slide("Priority Analysis", [
            f"Priority Score: {pa.get('priority_score', 'N/A')}",
            _truncate(company_explainer, 350),
        ])

        # ── SLIDE 4: Key Opportunity Drivers ─────────────────────────────────
        drivers_text = str(pa.get('key_opportunity_drivers', '') or '')
        driver_bullets = [b.strip('•- ').strip() for b in drivers_text.split('\n\n') if b.strip()][:6]
        if not driver_bullets:
            driver_bullets = [_truncate(drivers_text, 350)] if drivers_text else ["No opportunity driver data available."]
        _add_content_slide("Key Opportunity Drivers", driver_bullets)

        # ── SLIDE 5: Engagement Recommendation ───────────────────────────────
        eng_text = str(pa.get('engagement_recommendation', '') or '')
        eng_bullets = [b.strip('•- ').strip() for b in eng_text.split('\n\n') if b.strip()][:6]
        if not eng_bullets:
            eng_bullets = [_truncate(eng_text, 350)] if eng_text else ["No engagement recommendation data available."]
        _add_content_slide("Engagement Recommendation", eng_bullets)

        # ── SLIDE 6: Installed Base Overview ─────────────────────────────────
        inst_dd = customer_data.get('installed_base', []) or []
        if inst_dd:
            from collections import Counter
            eq_c = Counter(str(r.get('equipment_type') or r.get('equipment') or 'Unknown') for r in inst_dd)
            country_c = Counter(str(r.get('country_internal') or r.get('country') or 'Unknown') for r in inst_dd)
            years_raw = []
            for r in inst_dd:
                y = r.get('start_year_internal') or r.get('start_year') or r.get('year')
                try:
                    if y: years_raw.append(int(float(y)))
                except Exception: pass
            age_range = f"{min(years_raw)} – {max(years_raw)}" if years_raw else "N/A"
            _add_content_slide("Installed Base Overview", [
                f"Total Equipment Records: {len(inst_dd)}",
                f"Top Equipment Types: {', '.join(f'{k} ({v})' for k, v in eq_c.most_common(4))}",
                f"Countries Represented: {', '.join(f'{k} ({v})' for k, v in country_c.most_common(4))}",
                f"Startup Year Range: {age_range}",
                f"OEM Mix (top): {', '.join(str(r.get('oem') or r.get('manufacturer') or '') for r in inst_dd[:3] if r.get('oem') or r.get('manufacturer'))}",
            ])
        else:
            _add_content_slide("Installed Base Overview", ["No installed base data available."])

        # ── SLIDE 7: Order Intake History (table) ────────────────────────────
        if crm_history and isinstance(crm_history, dict):
            yearly_df = crm_history.get('yearly_df')
            if yearly_df is not None and not yearly_df.empty:
                _add_table_slide(
                    "Order Intake History",
                    ["Year", "Projects", "Total Value (EUR)", "Won Value (EUR)", "Win Rate %"],
                    [
                        [
                            str(rec.get('Year', '')),
                            str(int(rec.get('Projects', 0) or 0)),
                            f"{float(rec.get('Total Value (EUR)', 0) or 0):,.0f}",
                            f"{float(rec.get('Won Value (EUR)', 0) or 0):,.0f}",
                            f"{float(rec.get('Win Rate %', 0) or 0):.1f}%",
                        ]
                        for _, rec in yearly_df.head(12).iterrows()
                    ]
                )

        # ── SLIDE 8: Market Intelligence ─────────────────────────────────────
        mi = profile_data.get('market_intelligence', {})
        if mi:
            _add_content_slide("Market Intelligence", [
                _truncate(mi.get('financial_health', ''), 380),
                _truncate(mi.get('market_position', ''), 280),
                _truncate(mi.get('recent_developments', ''), 200),
                _truncate(mi.get('strategic_outlook', ''), 200),
                _truncate(mi.get('risk_assessment', ''), 200),
            ])

        # ── SLIDE 9: Country Intelligence ────────────────────────────────────
        ci = profile_data.get('country_intelligence', {})
        if ci:
            _add_content_slide("Country Intelligence", [
                _truncate(ci.get('steel_market_summary', ''), 280),
                _truncate(ci.get('economic_context', ''), 200),
                _truncate(ci.get('trade_tariff_context', ''), 200),
                _truncate(ci.get('investment_drivers', ''), 200),
            ])

        # ── SLIDE 10: Metallurgical Insights ─────────────────────────────────
        meta = profile_data.get('metallurgical_insights', {})
        if meta:
            _add_content_slide("Metallurgical & Technical Insights", [
                _truncate(meta.get('process_efficiency', ''), 280),
                _truncate(meta.get('carbon_footprint_strategy', ''), 200),
                _truncate(meta.get('modernization_potential', ''), 280),
                _truncate(meta.get('technical_bottlenecks', ''), 200),
            ])

        # ── SLIDES 11+: AI-generated narrative slides ─────────────────────────
        for item in (ai_slide_outline or [])[:6]:
            title_sl = str(item.get('title', 'Analysis'))
            bullets_sl = item.get('bullets', []) if isinstance(item.get('bullets'), list) else []
            if bullets_sl:
                _add_content_slide(title_sl, [str(b) for b in bullets_sl])

        # ── SLIDE: Value Proposition ──────────────────────────────────────────
        strat = profile_data.get('sales_strategy', {})
        vp = str(strat.get('value_proposition', '') or '')
        if vp:
            vp_bullets = [b.strip('•- ').strip() for b in vp.split('\n\n') if b.strip()][:6]
            if not vp_bullets:
                vp_bullets = [_truncate(vp, 350)]
            _add_content_slide("Value Proposition & Sales Strategy", vp_bullets)

        if strat.get('competitive_landscape'):
            _add_content_slide("Competitive Landscape", [
                _truncate(strat.get('competitive_landscape', ''), 380),
            ])

        # ── SLIDE: Recommended Next Steps ─────────────────────────────────────
        next_steps_raw = str(strat.get('suggested_next_steps', '') or '').split('\n')
        clean_steps = [s.strip('- ').strip() for s in next_steps_raw if s.strip()][:5]
        if not clean_steps:
            clean_steps = [
                "Schedule plant-level diagnostic visit to quantify actual bottlenecks.",
                "Prepare KPI-linked value case: yield, energy, uptime, CO2.",
                "Introduce SMS decarbonization roadmap in executive-level meeting.",
                "Sequence service quick wins, digital revamp, and capex modernization phases.",
                "Align technical champion workshop with customer process engineering team.",
            ]
        _add_content_slide("SMS Group — Recommended Next Steps", clean_steps)

        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
    
    def generate_comprehensive_pdf(
        self,
        customer_name: str,
        profile_data: Dict,
        customer_data: Dict,
        market_intel: Dict = None,
        projects: List[Dict] = None,
        financial_data: Dict = None,
        charts: Dict = None,
        crm_history: Dict = None,
        ib_data: Dict = None,
    ) -> BytesIO:
        """
        Generate comprehensive customer analysis PDF (15 sections) mirroring the DOCX report.
        """
        # Auto-populate market_intel from profile_data when not explicitly provided
        if not market_intel and profile_data:
            market_intel = profile_data.get('market_intelligence', {}) or {}

        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.lib import colors
            from reportlab.platypus import (
                SimpleDocTemplate, Paragraph, Spacer, PageBreak,
                Table, TableStyle, Image as RLImage, HRFlowable
            )
            from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
        except ImportError:
            logger.error("reportlab not available for PDF generation")
            raise ImportError("reportlab is required for PDF export. Install with: pip install reportlab")
        
        buffer = BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=A4,
            rightMargin=60,
            leftMargin=60,
            topMargin=60,
            bottomMargin=40,
        )

        elements = []
        styles = getSampleStyleSheet()

        # ── Style definitions ───────────────────────────────────────────────
        styles.add(ParagraphStyle(
            name='RptTitle', parent=styles['Heading1'],
            fontSize=26, textColor=colors.HexColor('#003366'),
            spaceAfter=20, alignment=TA_CENTER
        ))
        styles.add(ParagraphStyle(
            name='RptSubtitle', parent=styles['Heading2'],
            fontSize=18, textColor=colors.HexColor('#0066cc'),
            spaceAfter=14, alignment=TA_CENTER
        ))
        styles.add(ParagraphStyle(
            name='Rpt1', parent=styles['Heading1'],
            fontSize=16, textColor=colors.HexColor('#003366'),
            spaceAfter=10, spaceBefore=14,
        ))
        styles.add(ParagraphStyle(
            name='Rpt2', parent=styles['Heading2'],
            fontSize=13, textColor=colors.HexColor('#0066cc'),
            spaceAfter=8, spaceBefore=10,
        ))
        styles.add(ParagraphStyle(
            name='Rpt3', parent=styles['Heading3'],
            fontSize=11, textColor=colors.HexColor('#374151'),
            spaceAfter=6, spaceBefore=8, fontName='Helvetica-BoldOblique'
        ))
        styles.add(ParagraphStyle(
            name='RptBody', parent=styles['BodyText'],
            fontSize=10, textColor=colors.HexColor('#374151'),
            spaceAfter=8, leading=15
        ))
        styles.add(ParagraphStyle(
            name='RptBullet', parent=styles['BodyText'],
            fontSize=10, textColor=colors.HexColor('#374151'),
            spaceAfter=4, leftIndent=18, bulletIndent=6, leading=14
        ))
        styles.add(ParagraphStyle(
            name='RptSmall', parent=styles['BodyText'],
            fontSize=9, textColor=colors.HexColor('#6b7280'),
            spaceAfter=4, fontName='Helvetica-Oblique'
        ))

        seen_blocks = set()

        def _s(val):
            """Safe HTML-escaped string for ReportLab paragraphs."""
            if val is None:
                return 'N/A'
            return str(val).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br/>')

        def _norm(val):
            return ' '.join(str(val or '').strip().lower().split())

        def section_heading(number: str, title: str):
            elements.append(PageBreak())
            elements.append(Paragraph(f"<b>{number}. {title}</b>", styles['Rpt1']))
            elements.append(HRFlowable(width='100%', thickness=1.5, color=colors.HexColor('#003366')))
            elements.append(Spacer(1, 0.1 * inch))

        def sub_heading(title: str):
            elements.append(Spacer(1, 0.08 * inch))
            elements.append(Paragraph(f"<b>{title}</b>", styles['Rpt2']))

        import re as _re

        def body_para(text: str):
            """Render a text block as paragraph(s). Detects numbered/bulleted lists and renders
            each item as a separate bullet paragraph instead of one long run."""
            if not text or not str(text).strip() or str(text).strip() in ('N/A', 'None', '{}'):
                return
            raw = str(text).strip()
            normalized = _norm(raw)
            if not normalized or normalized in seen_blocks:
                return
            seen_blocks.add(normalized)
            lines = raw.split('\n')
            list_pat = _re.compile(r'^\s*(\d+[.)]\s+|[•\-\*]\s+)')
            list_lines = [l for l in lines if list_pat.match(l)]
            if len(list_lines) >= 2:
                current_para_lines = []
                for line in lines:
                    line_s = line.strip()
                    if not line_s:
                        if current_para_lines:
                            elements.append(Paragraph(_s(' '.join(current_para_lines)), styles['RptBody']))
                            current_para_lines = []
                    elif list_pat.match(line_s):
                        if current_para_lines:
                            elements.append(Paragraph(_s(' '.join(current_para_lines)), styles['RptBody']))
                            current_para_lines = []
                        clean = list_pat.sub('', line_s).strip()
                        elements.append(Paragraph(f"\u2022 {_s(clean)}", styles['RptBullet']))
                    else:
                        current_para_lines.append(line_s)
                if current_para_lines:
                    elements.append(Paragraph(_s(' '.join(current_para_lines)), styles['RptBody']))
            else:
                elements.append(Paragraph(_s(raw), styles['RptBody']))

        def kv_table(rows_data: list, col_widths=None):
            """Render a two-column key-value table."""
            if not rows_data:
                return
            cw = col_widths or [2.2 * inch, 4.4 * inch]
            t = Table(rows_data, colWidths=cw)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#eef2ff')),
                ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 10),
                ('GRID', (0, 0), (-1, -1), 0.4, colors.HexColor('#d1d5db')),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('ROWBACKGROUNDS', (0, 0), (-1, -1), [colors.white, colors.HexColor('#f9fafb')]),
            ]))
            elements.append(t)
            elements.append(Spacer(1, 0.12 * inch))

        def data_table(header_row: list, data_rows: list, col_widths=None):
            """Render a multi-column data table with a header row."""
            if not data_rows:
                return
            all_rows = [header_row] + data_rows
            n_cols = len(header_row)
            cw = col_widths or [6.6 * inch / n_cols] * n_cols
            t = Table(all_rows, colWidths=cw, repeatRows=1)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('GRID', (0, 0), (-1, -1), 0.4, colors.HexColor('#d1d5db')),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f9fafb')]),
                ('WORDWRAP', (0, 0), (-1, -1), True),
            ]))
            elements.append(t)
            elements.append(Spacer(1, 0.12 * inch))

        def add_chart_image(fig, caption: str):
            if not PLOTLY_AVAILABLE or fig is None:
                return
            try:
                img_bytes = pio.to_image(fig, format='png', width=760, height=380)
                elements.append(RLImage(BytesIO(img_bytes), width=6.0 * inch, height=3.0 * inch))
                elements.append(Paragraph(f"<i>{_s(caption)}</i>", styles['RptSmall']))
                elements.append(Spacer(1, 0.12 * inch))
            except Exception as e:
                logger.warning(f"Could not render chart '{caption}': {e}")

        # ── TITLE PAGE ──────────────────────────────────────────────────────
        logo_path = Path(__file__).resolve().parent.parent.parent / "assets" / "logo.png"
        if logo_path.exists():
            logo = RLImage(str(logo_path), width=1.6 * inch, height=0.45 * inch)
            logo.hAlign = 'RIGHT'
            elements.append(logo)

        elements.append(Spacer(1, 1.4 * inch))
        elements.append(Paragraph("<b>Customer Analysis Report</b>", styles['RptTitle']))
        elements.append(Spacer(1, 0.3 * inch))
        elements.append(Paragraph(f"<b>{_s(customer_name)}</b>", styles['RptSubtitle']))
        elements.append(Spacer(1, 0.25 * inch))
        elements.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')} | CONFIDENTIAL", styles['RptSmall']))
        elements[-1].alignment = TA_CENTER
        elements.append(Spacer(1, 0.6 * inch))
        elements.append(HRFlowable(width='80%', thickness=1, color=colors.HexColor('#003366')))

        # ── TABLE OF CONTENTS ────────────────────────────────────────────────
        elements.append(PageBreak())
        elements.append(Paragraph("<b>Table of Contents</b>", styles['Rpt1']))
        toc_items = [
            "1. Customer Profile",
            "2. Priority Ranking Analysis",
            "3. Project History & Sales Relationship",
            "4. Customer Interactions",
            "5. CRM Historical Performance",
            "6. Deep Dive Analytics",
            "7. Market Intelligence",
            "8. Country-Level Intelligence",
            "9. Installed Base Summary",
            "10. Project Analysis",
            "11. Metallurgical & Technical Insights",
            "12. Strategic Sales Pitch",
            "13. Financial Analysis",
            "14. Recent News & Developments",
            "15. References",
        ]
        for item in toc_items:
            elements.append(Paragraph(f"• {item}", styles['RptBullet']))

        # ── SECTION 1: CUSTOMER PROFILE ─────────────────────────────────────
        section_heading("1", "Customer Profile")
        sub_heading("1.1 Locations on Map")
        if charts and 'Locations Map' in charts:
            add_chart_image(charts['Locations Map'], "Global Footprint — Plant Locations")
        else:
            body_para("No geographic data available for map rendering.")

        sub_heading("1.2 Basic Information")
        basic = profile_data.get('basic_data', {})
        kv_table([
            ["Company Name",    _s(basic.get('name'))],
            ["Headquarters",    _s(basic.get('hq_address'))],
            ["CEO / Management",_s(basic.get('ceo') or basic.get('management'))],
            ["Owner / Parent",  _s(basic.get('owner'))],
            ["Employees (FTE)", _s(basic.get('fte'))],
            ["Industry Focus",  _s(basic.get('company_focus'))],
            ["Financial Status",_s(basic.get('financials'))],
            ["Ownership History",_s(basic.get('ownership_history'))],
        ])
        for title, key in [
            ("1.2.1 Corporate History & Origin", 'corporate_history'),
            ("1.2.2 Shareholding / Capital Structure", 'capital_structure'),
            ("1.2.3 Employees Breakdown", 'employee_breakdown'),
            ("1.2.4 Executive Board & Management Organization", 'executive_board'),
            ("1.2.5 Group Subsidiaries / Affiliates Overview", 'subsidiaries'),
        ]:
            val = basic.get(key, '')
            if val:
                sub_heading(title)
                body_para(val)

        # Locations detail table
        # Locations detail table — filtered to customer's primary country
        locs = profile_data.get('locations', [])
        if locs:
            sub_heading("1.3 Plant Locations")
            filter_country = (customer_data or {}).get('crm_country', '')
            if filter_country:
                filtered_locs = [loc for loc in locs if str(loc.get('country', '')).strip().lower() == filter_country.lower()]
                # Fallback: if no match, show all (cross-border HQ or data gap)
                filtered_locs = filtered_locs if filtered_locs else locs[:20]
            else:
                filtered_locs = locs[:20]
            loc_rows = []
            for loc in filtered_locs:
                loc_rows.append([
                    _s(loc.get('city', '')),
                    _s(loc.get('country', '')),
                    _s(loc.get('final_products', '')),
                    _s(loc.get('tons_per_year', '')),
                ])
            data_table(
                ["City", "Country", "Products", "Capacity (t/y)"],
                loc_rows,
                col_widths=[1.4*inch, 1.2*inch, 2.2*inch, 1.8*inch]
            )

        sub_heading("1.4 Equipment Distribution & Fleet Statistics")
        for chart_key in ['Equipment Distribution', 'Status Distribution', 'Age Distribution', 'Capacity Profile']:
            if charts and chart_key in charts:
                add_chart_image(charts[chart_key], chart_key)

        sub_heading("1.5 Statistical Data Analysis")
        stat = profile_data.get('statistical_interpretations', {})
        body_para(stat.get('charts_explanation') or "Installed-base distribution analysis based on available data.")

        # ── SECTION 2: PRIORITY RANKING ──────────────────────────────────────
        section_heading("2", "Priority Ranking Analysis")
        pa = profile_data.get('priority_analysis', {})
        kv_table([
            ["Priority Score", _s(pa.get('priority_score'))],
            ["Priority Rank",  _s(pa.get('priority_rank'))],
        ])
        explainer = pa.get('company_explainer') or pa.get('reasoning', '')
        if explainer:
            sub_heading("Company Analysis")
            body_para(explainer)
        if pa.get('key_opportunity_drivers'):
            sub_heading("Key Opportunity Drivers")
            body_para(pa['key_opportunity_drivers'])
        if pa.get('engagement_recommendation'):
            sub_heading("Engagement Recommendation")
            body_para(pa['engagement_recommendation'])

        # ── SECTION 3: PROJECT HISTORY & SALES RELATIONSHIP ─────────────────
        section_heading("3", "Project History & Sales Relationship")
        history = profile_data.get('history', {})
        context_pd = profile_data.get('context', {})
        kv_table([
            ["CRM Rating",          _s(history.get('crm_rating'))],
            ["SMS Relationship",    _s(history.get('sms_relationship'))],
            ["Key Contact Person",  _s(history.get('key_person'))],
            ["Latest Visits",       _s(history.get('latest_visits'))],
            ["Total Won Value (EUR)",_s(history.get('total_won_value_eur'))],
            ["Win Rate (%)",        _s(history.get('win_rate_pct'))],
        ])
        for key, label in [('latest_projects', 'Latest Projects'), ('realized_projects', 'Realized Projects')]:
            val = history.get(key, '')
            if val:
                sub_heading(label)
                body_para(val)
        for key, label in [
            ('sms_delivery_history', 'Existing Facilities of SMS Group (History)'),
            ('current_projects_detail', 'Current Projects (Active Opportunities)'),
            ('projects_under_execution', 'Projects Under Execution'),
            ('lost_projects', 'Lost Projects'),
        ]:
            val = history.get(key, '')
            if val:
                sub_heading(label)
                body_para(val)
        announced = (profile_data.get('market_intelligence') or {}).get('announced_investments', '')
        if announced:
            sub_heading('Announced Investments (News)')
            body_para(announced)
        if context_pd.get('end_customer'):
            sub_heading("End Customer / Supply Chain")
            body_para(context_pd['end_customer'])
        if context_pd.get('market_position'):
            sub_heading("Market Position")
            body_para(context_pd['market_position'])

        # ── SECTION 4: CUSTOMER INTERACTIONS ─────────────────────────────────
        section_heading("4", "Customer Interactions")
        interaction_summary = profile_data.get('customer_interaction_summary', {}) or {}
        interactions = profile_data.get('customer_interactions', []) or []
        if interaction_summary:
            kv_table([
                ["Total Interactions", _s(interaction_summary.get('total_interactions', 0))],
                ["Last Contact Date", _s(interaction_summary.get('last_contact_date'))],
                ["Last Contact Location", _s(interaction_summary.get('last_contact_location'))],
                ["Last Contact Owner", _s(interaction_summary.get('last_contact_owner'))],
                ["Last Contact Subject", _s(interaction_summary.get('last_contact_subject'))],
                ["Main Channels", _s(', '.join(interaction_summary.get('top_channels', [])) or 'N/A')],
                ["Main SMS Contacts", _s(', '.join(interaction_summary.get('top_contacts', [])) or 'N/A')],
            ])
        if interactions:
            rows = []
            for item in interactions[:12]:
                rows.append([
                    _s(str(item.get('start_dt', '') or '')[:10]),
                    _s(item.get('account', '')),
                    _s(' | '.join([x for x in [item.get('distribution_channel', ''), item.get('meeting_location', '')] if x])),
                    _s(item.get('employee_responsible', '')),
                    _s(item.get('subject', '')),
                ])
            sub_heading("Recent Interaction Timeline")
            data_table(
                ["Date", "Account", "Channel / Location", "Responsible", "Subject"],
                rows,
                col_widths=[0.9*inch, 1.5*inch, 1.6*inch, 1.2*inch, 1.4*inch],
            )
        elif not interaction_summary:
            body_para("No recent customer interactions available in the SAP Sales Cloud visit export.")

        # ── SECTION 5: CRM HISTORICAL PERFORMANCE ────────────────────────────
        section_heading("5", "CRM Historical Performance")
        if crm_history:
            src = crm_history.get('source', '')
            if src:
                elements.append(Paragraph(f"<i>Data source: {_s(src)}</i>", styles['RptSmall']))
            yearly_df = crm_history.get('yearly_df')
            if yearly_df is not None and not yearly_df.empty:
                sub_heading("Year-by-Year Summary")
                rows = []
                for _, r in yearly_df.iterrows():
                    # Force year to integer string (avoid "2024.0")
                    raw_year = r.get('Year', '')
                    try:
                        year_str = str(int(float(raw_year))) if raw_year != '' and raw_year is not None else ''
                    except (ValueError, TypeError):
                        year_str = str(raw_year)
                    rows.append([
                        year_str,
                        str(int(r.get('Projects', 0) or 0)),
                        f"EUR {float(r.get('Total Value (EUR)', 0) or 0):,.0f}",
                        f"EUR {float(r.get('Won Value (EUR)', 0) or 0):,.0f}",
                        f"{float(r.get('Win Rate %', 0) or 0):.1f}%",
                    ])
                data_table(
                    ["Year", "Projects", "Total Pipeline (EUR)", "Won Value (EUR)", "Win Rate %"],
                    rows,
                    col_widths=[0.8*inch, 0.9*inch, 1.8*inch, 1.8*inch, 1.3*inch]
                )
            won_list = crm_history.get('won_list')
            if won_list is not None and not won_list.empty:
                sub_heading("Won Projects Detail")
                show_cols = [c for c in [
                    'account_name', 'codeword_sales', 'customer_project',
                    'cp_expected_value_eur', '_year', 'account_country', 'sp_coe',
                ] if c in won_list.columns]
                if show_cols:
                    won_sub = won_list[show_cols].head(20)
                    hdr = [c.replace('_', ' ').title() for c in show_cols]
                    d_rows = []
                    for _, dr in won_sub.iterrows():
                        d_rows.append([_s(dr.get(c, '')) for c in show_cols])
                    data_table(hdr, d_rows)
        else:
            body_para("No CRM historical data available for this customer.")

        # ── SECTION 6: DEEP DIVE ANALYTICS ───────────────────────────────────
        section_heading("6", "Deep Dive Analytics")
        proj_dd = customer_data.get('projects', []) or []
        inst_dd = customer_data.get('installed_base', []) or []
        total_rev = sum(p.get('value', 0) or 0 for p in proj_dd)
        active_p = sum(1 for p in proj_dd if p.get('status') in ('Active', 'In Progress'))
        kv_table([
            ["Total Projects",   str(len(proj_dd))],
            ["Active Projects",  str(active_p)],
            ["Total Equipment (CRM Export)",  str(len(inst_dd))],
            ["Total CRM Revenue",f"EUR {total_rev:,.0f}"],
        ])
        if inst_dd:
            from collections import Counter
            eq_counter = Counter(str(r.get('equipment_type') or r.get('equipment') or 'Unknown') for r in inst_dd)
            sub_heading("Equipment Family Breakdown")
            eq_rows = [[eq, str(cnt)] for eq, cnt in eq_counter.most_common(20)]
            data_table(["Equipment Type", "Count"], eq_rows, col_widths=[4.5*inch, 2.1*inch])

        # ── SECTION 7: MARKET INTELLIGENCE ───────────────────────────────────
        section_heading("7", "Market Intelligence")
        if market_intel:
            for title, key in [
                ("Financial Health",    'financial_health'),
                ("Recent Developments", 'recent_developments'),
                ("Market Position",     'market_position'),
                ("Strategic Outlook",   'strategic_outlook'),
                ("Risk Assessment",     'risk_assessment'),
                ("Market Size",         'market_size'),
                ("Growth Trends",       'growth_trends'),
            ]:
                val = market_intel.get(key, '')
                if isinstance(val, dict):
                    val = val.get('summary', '') or val.get('text', '') or str(val)
                if val:
                    sub_heading(title)
                    body_para(val)
            if market_intel.get('competitors'):
                sub_heading("Key Competitors")
                for comp in market_intel['competitors']:
                    elements.append(Paragraph(f"• {_s(comp)}", styles['RptBullet']))
            for title, key in [
                ("Workforce Strategy", 'workforce_strategy'),
                ("Products of Steel Business Unit", 'product_portfolio'),
                ("Construction / Automotive-specific Market Sections", 'end_market_breakdown'),
            ]:
                val = market_intel.get(key, '')
                if val:
                    sub_heading(title)
                    body_para(val)
        else:
            body_para("Market intelligence data not available. Generate or refresh the customer profile to populate this section.")

        # ── SECTION 8: COUNTRY-LEVEL INTELLIGENCE ────────────────────────────
        section_heading("8", "Country-Level Intelligence")
        ci = profile_data.get('country_intelligence', {})
        if ci:
            for title, key in [
                ("Steel Market Summary",        'steel_market_summary'),
                ("Economic Context",            'economic_context'),
                ("Trade & Tariff Context",      'trade_tariff_context'),
                ("Automotive Sector Trends",    'automotive_sector'),
                ("Investment Drivers",          'investment_drivers'),
            ]:
                val = ci.get(key, '')
                if val and val not in ('N/A', 'None'):
                    sub_heading(title)
                    body_para(val)
        else:
            body_para("Country intelligence data not available.")

        # ── SECTION 9: INSTALLED BASE SUMMARY ────────────────────────────────
        section_heading("9", "Installed Base Summary")
        if ib_data and ib_data.get('n_units', 0) > 0:
            kv_table([
                ["Equipment Units (IB List)",      str(ib_data.get('n_units', 0))],
                ["Average Age (years)",  str(ib_data.get('avg_age', 'N/A'))],
                ["Equipment Types",      ', '.join(str(t) for t in ib_data.get('equipment_types', [])) or 'N/A'],
                ["Countries / Regions",  ', '.join(str(c) for c in ib_data.get('countries', [])) or 'N/A'],
            ])
            df_ib = ib_data.get('df')
            if df_ib is not None and not df_ib.empty:
                sub_heading("Equipment Records (top 30)")
                display_cols = [c for c in [
                    'ib_machine', 'ib_description', 'ib_product',
                    'ib_city', 'ib_customer_country', 'ib_startup', '_age', 'ib_status',
                ] if c in df_ib.columns]
                if display_cols:
                    subset = df_ib[display_cols].head(30)
                    hdr = [c.replace('ib_', '').replace('_', ' ').title() for c in display_cols]
                    d_rows = []
                    for _, dr in subset.iterrows():
                        val = dr.get(display_cols[0], '')
                        d_rows.append([_s(dr.get(c, '') if str(dr.get(c, '')) != 'nan' else '') for c in display_cols])
                    data_table(hdr, d_rows)
        else:
            body_para("No installed base records found for this customer.")

        # ── SECTION 10: PROJECT ANALYSIS ──────────────────────────────────────
        section_heading("10", "Project Analysis")
        if proj_dd:
            active_c = sum(1 for p in proj_dd if p.get('status') in ('Active', 'In Progress'))
            won_c = sum(1 for p in proj_dd if p.get('status') in ('Completed', 'Won'))
            kv_table([
                ["Total Projects",     str(len(proj_dd))],
                ["Active / In Progress", str(active_c)],
                ["Completed / Won",    str(won_c)],
                ["Total Value",        f"EUR {total_rev:,.0f}"],
            ])
            sub_heading("Project Details")
            proj_rows = []
            for p in proj_dd[:30]:
                proj_rows.append([
                    _s(p.get('name', 'Unnamed')),
                    _s(p.get('status', 'N/A')),
                    f"EUR {p.get('value', 0) or 0:,.0f}",
                    _s(p.get('start_date', 'N/A')),
                    _s(p.get('end_date', 'N/A')),
                ])
            data_table(
                ["Project Name", "Status", "Value (EUR)", "Start", "End"],
                proj_rows,
                col_widths=[2.4*inch, 1.0*inch, 1.2*inch, 1.0*inch, 1.0*inch]
            )
        else:
            body_para("No detailed project records available for this customer.")

        # ── SECTION 11: METALLURGICAL INSIGHTS ───────────────────────────────
        section_heading("11", "Metallurgical & Technical Insights")
        meta = profile_data.get('metallurgical_insights', {})
        if meta:
            for label, key in [
                ("Process Efficiency",            'process_efficiency'),
                ("Carbon Footprint / Green Steel", 'carbon_footprint_strategy'),
                ("Modernization Potential",        'modernization_potential'),
                ("Technical Bottlenecks",          'technical_bottlenecks'),
            ]:
                val = meta.get(key, '')
                if val:
                    sub_heading(label)
                    body_para(val)
        else:
            body_para("No metallurgical insights available. Generate a profile first.")

        # ── SECTION 12: STRATEGIC SALES PITCH ────────────────────────────────
        section_heading("12", "Strategic Sales Pitch")
        strat = profile_data.get('sales_strategy', {})
        if strat:
            for label, key in [
                ("Commercial Sales SMS group - Sales Structure", 'sms_commercial_structure'),
                ("Buying Center Map", 'buying_center_map'),
            ]:
                val = strat.get(key, '')
                if val:
                    sub_heading(label)
                    body_para(val)
            for label, key in [
                ("Recommended Portfolio",  'recommended_portfolio'),
                ("Value Proposition",      'value_proposition'),
                ("Competitive Landscape",  'competitive_landscape'),
                ("Suggested Next Steps",   'suggested_next_steps'),
            ]:
                val = strat.get(key, '')
                if val:
                    sub_heading(label)
                    body_para(val)
            if strat.get('compliance_guidance'):
                sub_heading('Compliance Implications For SMS')
                body_para(strat.get('compliance_guidance'))
            if not any(strat.get(k) for k in ('recommended_portfolio', 'value_proposition', 'competitive_landscape', 'suggested_next_steps')):
                body_para(str(strat))
        else:
            body_para("No sales strategy available. Generate a profile first.")

        # ── SECTION 13: FINANCIAL ANALYSIS ───────────────────────────────────
        section_heading("13", "Financial Analysis")
        if financial_data:
            costs = financial_data.get('cost_breakdown', {})
            if costs:
                sub_heading("Cost Breakdown")
                cost_rows = [[str(k), f"EUR {float(v):,.2f}" if str(v).replace('.', '').isdigit() else str(v)] for k, v in costs.items()]
                data_table(["Category", "Amount (EUR)"], cost_rows, col_widths=[3.3*inch, 3.3*inch])
            variance = financial_data.get('budget_variance', {})
            if variance:
                sub_heading("Budget Variance Analysis")
                kv_table([
                    ["Budgeted",  f"EUR {variance.get('budgeted', 0):,.2f}"],
                    ["Actual",    f"EUR {variance.get('actual', 0):,.2f}"],
                    ["Variance",  f"EUR {variance.get('variance', 0):,.2f} ({variance.get('variance_percent', 0):.1f}%)"],
                    ["Status",    _s(variance.get('status', 'Unknown'))],
                ])
        else:
            body_para("Detailed financial data not separately loaded. See financial commentary in Sections 2 and 6.")

        # ── SECTION 14: RECENT NEWS & DEVELOPMENTS ───────────────────────────
        section_heading("14", "Recent News & Developments")
        news = profile_data.get('recent_news', [])
        if news:
            for n in news[:10]:
                title_txt = _s(n.get('title', 'Untitled'))
                date_txt  = _s(n.get('published_date', ''))
                src_txt   = _s(n.get('source', ''))
                elements.append(Paragraph(f"<b>{title_txt}</b>  ({date_txt}{' | ' + src_txt if src_txt else ''})", styles['RptBullet']))

                # 3-4 line summary per news item (roughly up to four sentences)
                desc_raw = str(n.get('description') or n.get('summary') or '').strip()
                if desc_raw:
                    sentences = [s.strip() for s in desc_raw.replace('\n', ' ').split('. ') if s.strip()]
                    summary = '. '.join(sentences[:4])
                    if summary and not summary.endswith('.'):
                        summary += '.'
                    if len(summary) > 650:
                        summary = summary[:647] + '...'
                    elements.append(Paragraph(_s(summary), ParagraphStyle(
                        name='NewsBody', parent=styles['RptBody'],
                        leftIndent=22, fontSize=9, spaceAfter=6
                    )))
        else:
            body_para("No recent news available.")

        # ── SECTION 15: REFERENCES ────────────────────────────────────────────
        section_heading("15", "References")
        refs = profile_data.get('references', []) or profile_data.get('source_links', [])
        if refs:
            for ref in refs:
                elements.append(Paragraph(f"• {_s(ref)}", styles['RptBullet']))
        else:
            body_para("No specific external references cited.")

        # ── BUILD ─────────────────────────────────────────────────────────────
        def _pdf_footer(canvas, doc):  # noqa: E306
            canvas.saveState()
            canvas.setFont('Helvetica', 8)
            canvas.setFillColor(colors.HexColor('#9ca3af'))
            page_w = canvas._pagesize[0]
            page_num = canvas.getPageNumber()
            text = f"Customer Analysis Report  |  CONFIDENTIAL  |  Page {page_num}  |  Generated {datetime.now().strftime('%Y-%m-%d')}"
            canvas.drawCentredString(page_w / 2.0, 0.28 * inch, text)
            canvas.restoreState()

        doc.build(elements, onFirstPage=_pdf_footer, onLaterPages=_pdf_footer)
        buffer.seek(0)
        return buffer


# Singleton instance
enhanced_export_service = EnhancedExportService()
